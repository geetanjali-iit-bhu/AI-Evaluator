"""
ingestion.py
Loads PDF, PPTX, DOCX, TXT → chunks → LangChain Documents
"""

import os
from pathlib import Path
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter


def _load_pdf(path: str) -> str:
    import pdfplumber
    pages = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            for table in page.extract_tables():
                for row in table:
                    text += "\n" + " | ".join(c.strip() for c in row if c)
            pages.append(f"[Page {i}]\n{text}")
    return "\n\n".join(pages)


def _load_pptx(path: str) -> str:
    from pptx import Presentation
    texts = []
    for i, slide in enumerate(Presentation(path).slides, 1):
        texts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    texts.append(" | ".join(
                        cell.text.strip() for cell in row.cells
                    ))
    return "\n".join(texts)


def _load_docx(path: str) -> str:
    import docx
    doc = docx.Document(path)
    texts = [para.text for para in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            texts.append(" | ".join(
                cell.text.strip() for cell in row.cells
            ))
    return "\n".join(texts)


def _load_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


LOADERS = {
    ".pdf":  _load_pdf,
    ".pptx": _load_pptx,
    ".docx": _load_docx,
    ".txt":  _load_txt,
}


def load_documents(source_dir: str = "data") -> list[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=100,
        separators=["\n\n", "\n", ".", " ", ""],
    )

    documents: list[Document] = []

    for file_path in Path(source_dir).rglob("*"):
        suffix = file_path.suffix.lower()
        if suffix not in LOADERS:
            continue

        print(f"  Loading: {file_path}")
        try:
            raw_text = LOADERS[suffix](str(file_path))
        except Exception as e:
            print(f"    ⚠ Skipped {file_path.name}: {e}")
            continue

        if not raw_text.strip():
            print(f"    ⚠ Empty content in {file_path.name}, skipping.")
            continue

        chunks = splitter.create_documents(
            texts=[raw_text],
            metadatas=[{"source": file_path.name, "type": suffix.lstrip(".")}],
        )
        documents.extend(chunks)
        print(f"    ✓ {len(chunks)} chunks from {file_path.name}")

    return documents


if __name__ == "__main__":
    docs = load_documents("data")
    print(f"\nTotal chunks loaded: {len(docs)}")
    if docs:
        print("\nSample chunk:\n", docs[0].page_content[:300])